""" Work Center Analysis
The goals of this program are to automate the creation of
charts that show production numbers and build a slideshow
with titles of equipment along with operator stats.
"""
from openpyxl import Workbook, load_workbook
from openpyxl.utils import get_column_letter
import types
import os
import time
import warnings
from matplotlib import pyplot as plt
from matplotlib import interactive
import numpy as np
import textwrap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
warnings.simplefilter("ignore")

#https://matplotlib.org/stable/api/_as_gen/matplotlib.pyplot.show.html tutorials
def main():
    
    def wrap_labels(ax, width, break_long_words=False):
        labels = []
        
        for label in ax.get_xticklabels():
            
            text = label.get_text()
            labels.append(textwrap.fill(text, width=width,
                      break_long_words=break_long_words))
            
        ax.set_xticklabels(labels, rotation=0)
        
    stats = []
    stats2 = []
    statlist = []
    df = []
    
    pddep = ''
    masterpath = 'M:/Schedule/KPI/Press Reports/'
    deps = open(r'M:\Schedule\KPI\Press Reports\New WCA/deps.txt','r')
    deps = deps.read()
    deps = deps.replace('\n','').split(',')
    opps = open(r'M:\Schedule\KPI\Press Reports\New WCA/opps.txt','r')
    opps = opps.read()
    opps = opps.replace('\n','').split('!')
    stats = open(r'M:\Schedule\KPI\Press Reports\New WCA/Averages.txt','r')
    stats = stats.read()
    stats = stats.replace('\n','').split(',')
    
    for i in stats:
        
        item = i.split('!')
        stats2.append(item)

    stats = stats2
    
    for i in stats:

        i[2] = int(i[2])
        i[3] = int(i[3])
        i[4] = int(i[4])
        i[5] = int(i[5])
        
    while True:
        
        month = input('Month:')
        day = input('Day:')
        year = input('Year:')
        path = masterpath+year+'/'+month+'-'+day+'-'+year+'.xlsx'
        exist = os.path.exists(path)
                
        if exist == True:
            break

        else:        
            print('That is not a valid file.')
            time.sleep(1)


    wb = load_workbook(path)
    ws = wb.active

    for i in range(1,ws.max_row+1):
        
        cell = ws.cell(row=i,column=1)
        
        if cell.value == 'Cletus Lee Landsteiner':
            cell.value = 'Lee Landsteiner'

    #recodes press names
    for i in range(1,ws.max_row+1):
        
        cell = ws.cell(row=i,column=2)

        if cell.value == 'Sanden 07 992-18" FixedSize':
            cell.value = 'Sanden 7'
                        
        elif cell.value == 'Sanden 03 893-27" 1200' or cell.value == 'Sanden 04  910-27" 1200':
            cell.value = 'Sanden 3&4'
                        
        elif cell.value == 'Apollo 07':
            cell.value = 'Apollo 7'
                        
        elif cell.value == 'Apollo 04 - Navitor' or cell.value == 'Apollo 04':
            cell.value = 'Apollo 4'
                        
        elif cell.value == 'Webcom 16' or cell.value == 'Webcom 17' or cell.value == 'Webcom 13':
            cell.value = 'Webcom 16&17'

    for i in range(1,ws.max_row+1):
        
        statsfound = 0        
        cell = ws.cell(row=i,column=2)

        if cell.value in deps:
            currentdep = cell.value

            for j in range(i,ws.max_row+1):

                cell = ws.cell(row=j,column=1)

                if statsfound == 1:
                    break

                elif cell.value in opps:
                    currentopp = cell.value

                    for k in range(j,ws.max_row+1):
                        
                        cell = ws.cell(row=k,column=1)      

                        if cell.value in opps and cell.value != currentopp:
                            break
                        
                        elif cell.value == 'Employee Total':
                            hours = ws.cell(row=k,column=2)
                            gross = ws.cell(row=k,column=3)
                            net = ws.cell(row=k,column=4)
                            waste = ws.cell(row=k,column=5)
                            stats.insert(100,[currentdep,currentopp,hours.value,gross.value,net.value,waste.value])

                        elif cell.value == 'Work Center Total':
                            statsfound = 1
                            
                        if statsfound == 1:
                            break


    dir1 = 'M:/Schedule/KPI/Press Reports/New WCA/Hour'
    dir2 = 'M:/Schedule/KPI/Press Reports/New WCA/Waste'
    
    for i in os.listdir(dir1):
        
        os.remove(os.path.join(dir1,i))

    for i in os.listdir(dir2):
        
        os.remove(os.path.join(dir2,i))
        
    for department in deps:

        goodfeet = []
        goodfeetperhour = []
        waste = []
        opplist = []

        plot = False
            
        for stat in stats:
            
            if stat[0] == department:
                goodfeet.insert(100,stat[4])
                waste.insert(100,stat[5])
                opplist.insert(100,stat[1])
                goodfeetperhour.insert(100,round(stat[4]/stat[2]))
                
        n = len(opplist)
        
        if n >= 2:
            plot = True
                
        if plot == True:
            arange = np.arange(n)
            fig, ax = plt.subplots(figsize=(7,7))

            bar_width = .5
            feet = ax.bar(arange,goodfeet,bar_width,label='Good Feet')
            badfeet = ax.bar(arange,waste,bar_width,label='Waste',bottom=goodfeet)
            
            ax.set_ylabel('Feet')
            ax.set_xticks(arange)
            ax.set_xticklabels(opplist,wrap=True)
            
            ax.margins(x=.125)

            ax.legend(loc='upper right')

            ax.bar_label(feet,label_type='center',color='black')
            ax.bar_label(badfeet,label_type='center',color='black')
            wrap_labels(ax, 10)
            ax.legend(loc='upper right',bbox_to_anchor=(1,1.1))

            ax.set_axisbelow(True)
            plt.grid(axis = 'y',color = 'gray')
            plt.title('Good Feet/Waste',loc='center',fontsize = 20)
            
            plt.savefig(r'M:\Schedule\KPI\Press Reports\New WCA\Waste/'+department+'.png',)  

            ########
            fig, ax = plt.subplots(figsize=(7,7))

            bar_width = .5
            feet = ax.bar(arange,goodfeetperhour,bar_width,label='Good Feet')
            
            ax.set_ylabel('Feet')
            ax.set_xticks(arange)
            ax.set_xticklabels(opplist,wrap=True)
            
            ax.margins(x=.125)

            ax.legend(loc='upper right')

            ax.bar_label(feet,label_type='edge',color='black')
            wrap_labels(ax, 10)
            ax.legend(loc='upper right',bbox_to_anchor=(1,1.075))

            ax.set_axisbelow(True)
            plt.grid(axis = 'y',color = 'gray')
            plt.title('Good Feet Per Hour',loc='center',fontsize = 20)
            
            plt.savefig(r'M:\Schedule\KPI\Press Reports\New WCA\Hour/'+department+'.png',)

    ##
    prs = Presentation(r'M:\Schedule\KPI\Press Reports\New WCA/MasterFile.pptx')
    layout = prs.slide_layouts[5]
    title = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title)
    title = slide.shapes.title
    title.text = 'Precision Press Specialty Productivity'+'\n'+month+'-'+day+'-'+year
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    sub = slide.placeholders[1]
    sub = sub.element
    sub.getparent().remove(sub)
    
    for department in deps:
        hours = 'M:/Schedule/KPI/Press Reports/New WCA/Hour/'+department+'.png'
        waste = 'M:/Schedule/KPI/Press Reports/New WCA/Waste/'+department+'.png'
        exist1 = os.path.exists(hours)
        
        if exist1 == True:  
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            title.text = department
            title.top = Inches(.5)
            title.left = Inches(0)
            title.height = Inches(1)
            title.width = Inches(13)
            title_para = slide.shapes.title.text_frame.paragraphs[0]
            title_para.font.size = Pt(60)
            title_para.alignment = PP_ALIGN.CENTER
            
            left = Inches(1.75)
            top = Inches(1.5)
            pic = slide.shapes.add_picture(hours,left,top,Inches(5))

            left = Inches(6.75)
            top = Inches(1.5)
            pic2 = slide.shapes.add_picture(waste,left,top,Inches(5))            

    prs.save(r'M:\Schedule\KPI\Press Reports\WCA Presentation.pptx')
    os.startfile(r'M:\Schedule\KPI\Press Reports\WCA Presentation.pptx')        
    print('Program completed...')
    time.sleep(1)
    
if __name__ == '__main__':
    main() 
